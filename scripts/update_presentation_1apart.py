#!/usr/bin/env python3
"""Обновление презентации 1apart: 8 → 12 слайдов по ТЗ."""
from __future__ import annotations

import shutil
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
PRESENTATIONS = ROOT / "docs" / "presentations"
ASSETS = PRESENTATIONS / "assets"

C = {
    "eyebrow": RGBColor(0xA9, 0x82, 0x4F),
    "title": RGBColor(0x1B, 0x1B, 0x1B),
    "body": RGBColor(0x1B, 0x1B, 0x1B),
    "accent": RGBColor(0x3B, 0x17, 0x0B),
    "muted": RGBColor(0x7A, 0x75, 0x6C),
    "green": RGBColor(0x2E, 0x7D, 0x32),
    "gold": RGBColor(0xDC, 0xB1, 0x79),
    "gold_light": RGBColor(0xC1, 0x9B, 0x6A),
    "white": RGBColor(0xFF, 0xFF, 0xFF),
    "bg": RGBColor(0xFA, 0xF8, 0xF5),
    "card": RGBColor(0xFF, 0xFF, 0xFF),
    "border": RGBColor(0xE8, 0xE0, 0xD5),
    "dark": RGBColor(0x12, 0x12, 0x12),
}

FOOTER = "Первый Апарт-отель · Система автоматической отчётности"


def find_source() -> Path:
    candidates = [
        PRESENTATIONS / "archive" / "Презентация_1apart_с_заметками.pptx [Repaired].pptx",
        PRESENTATIONS / "Презентация_1apart_с_заметками.pptx [Repaired].pptx",
    ]
    for path in candidates:
        if path.exists():
            return path
    ext = Path(r"c:\Users\user\Documents\Люди")
    return next(
        p for p in ext.rglob("*.pptx") if "1apart" in p.name.lower() and "repaired" in p.name.lower()
    )


def output_path() -> Path:
    return PRESENTATIONS / "Презентация_1apart_с_заметками_12слайдов.pptx"


def _set_run(run, *, size: float, bold: bool = False, color: RGBColor | None = None) -> None:
    run.font.size = Pt(size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color


def _fill_bg(slide, rgb: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def _add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text: str,
    *,
    size: float = 24,
    bold: bool = False,
    color: RGBColor | None = None,
    align=PP_ALIGN.LEFT,
) -> None:
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    _set_run(run, size=size, bold=bold, color=color or C["body"])


def _add_bullets(slide, left, top, width, height, lines: list[str], *, size: float = 18) -> None:
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = 0
        p.space_after = Pt(6)
        run = p.add_run()
        run.text = line
        _set_run(run, size=size, color=C["body"])


def _add_card(slide, left, top, width, height) -> None:
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = C["card"]
    shape.line.color.rgb = C["border"]


def _footer(slide, num: int) -> None:
    _add_textbox(slide, Inches(0.55), Inches(7.0), Inches(8), Inches(0.25), FOOTER, size=9, color=C["muted"])
    _add_textbox(slide, Inches(12.2), Inches(7.0), Inches(0.5), Inches(0.25), str(num), size=9, color=C["muted"], align=PP_ALIGN.RIGHT)


def _notes(slide, text: str) -> None:
    ns = slide.notes_slide
    try:
        ph = ns.notes_placeholder
        ph.text = text
        return
    except Exception:
        pass
    tf = ns.notes_text_frame
    if tf is not None:
        tf.clear()
        tf.text = text


def _replace_shape_text(slide, old_prefix: str, new_text: str) -> bool:
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text_frame.text.strip().startswith(old_prefix):
            shape.text_frame.clear()
            p = shape.text_frame.paragraphs[0]
            run = p.add_run()
            run.text = new_text
            return True
    return False


def _replace_exact(slide, old: str, new: str) -> bool:
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text_frame.text.strip() == old:
            shape.text_frame.text = new
            return True
    return False


def build_content_slide(
    prs: Presentation,
    *,
    eyebrow: str,
    title: str,
    bullets: list[str] | None = None,
    body: str = "",
    image: Path | None = None,
    image_right: bool = True,
    slide_num: int,
    notes: str,
    extra_blocks: list[tuple[str, str, RGBColor]] | None = None,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    _fill_bg(slide, C["bg"])
    _add_textbox(slide, Inches(0.55), Inches(0.45), Inches(10), Inches(0.35), eyebrow, size=12, bold=True, color=C["eyebrow"])
    _add_textbox(slide, Inches(0.55), Inches(0.85), Inches(11.5 if image else 12), Inches(0.7), title, size=26, bold=True, color=C["title"])

    text_left = Inches(0.55)
    text_width = Inches(6.0 if image else 12.0)
    top = Inches(1.65)
    if body:
        _add_textbox(slide, text_left, top, text_width, Inches(0.9), body, size=16, color=C["body"])
        top = Inches(2.35)
    if bullets:
        _add_bullets(slide, text_left, top, text_width, Inches(3.8), bullets, size=16)
    if extra_blocks:
        y = Inches(4.8)
        for head, txt, color in extra_blocks:
            _add_textbox(slide, text_left, y, text_width, Inches(0.35), head, size=18, bold=True, color=color)
            _add_textbox(slide, text_left, y + Inches(0.35), text_width, Inches(0.55), txt, size=15, color=C["body"])
            y += Inches(1.0)
    if image and image.exists():
        ix = Inches(7.0) if image_right else Inches(0.55)
        slide.shapes.add_picture(str(image), ix, Inches(1.55), width=Inches(5.8))
    _footer(slide, slide_num)
    _notes(slide, notes)


def build_reco_slide(prs: Presentation, slide_num: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    _fill_bg(slide, C["bg"])
    _add_textbox(slide, Inches(0.55), Inches(0.45), Inches(10), Inches(0.35), "ВНЕДРЕНИЕ", size=12, bold=True, color=C["eyebrow"])
    _add_textbox(
        slide,
        Inches(0.55),
        Inches(0.85),
        Inches(12),
        Inches(0.9),
        "РЕКОМЕНДАЦИЯ — ЭТО ГОТОВАЯ ИНСТРУКЦИЯ ДЛЯ МЕНЕДЖЕРА",
        size=24,
        bold=True,
        color=C["title"],
    )
    img = ASSETS / "reco_flow.png"
    if img.exists():
        slide.shapes.add_picture(str(img), Inches(0.55), Inches(1.7), width=Inches(4.0))
    _add_textbox(slide, Inches(4.9), Inches(1.75), Inches(7.5), Inches(0.45), "Каждая рекомендация отвечает на четыре вопроса:", size=16, bold=True, color=C["accent"])
    qs = [
        "1. Что происходит и почему это важно?",
        "2. Что именно нужно сделать?",
        "3. Как проверить, что действие сработало?",
        "4. Что делать, если результат не достигнут?",
    ]
    _add_bullets(slide, Inches(4.9), Inches(2.25), Inches(7.5), Inches(2.0), qs, size=15)
    _add_card(slide, Inches(0.55), Inches(5.05), Inches(12.0), Inches(1.35))
    _add_textbox(slide, Inches(0.8), Inches(5.2), Inches(11.5), Inches(0.3), "💡 Карточку можно:", size=16, bold=True, color=C["accent"])
    extras = [
        "• принять или отложить;",
        "• назначить сотруднику;",
        "• отметить как выполненную;",
        "• выгрузить в Word для работы.",
    ]
    _add_bullets(slide, Inches(0.8), Inches(5.55), Inches(11.5), Inches(0.8), extras, size=15)
    _footer(slide, slide_num)
    _notes(
        slide,
        "Рекомендация — не просто цифра, а готовая инструкция: что происходит, что сделать, "
        "как проверить результат и что делать при отклонении. Карточку можно принять, отложить, "
        "назначить сотруднику и выгрузить в Word. Автоматической смены цен в TravelLine нет.",
    )


def build_perspectives_slide(prs: Presentation, slide_num: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    _fill_bg(slide, C["bg"])
    _add_textbox(slide, Inches(0.55), Inches(0.45), Inches(10), Inches(0.35), "РАЗВИТИЕ", size=12, bold=True, color=C["eyebrow"])
    _add_textbox(
        slide,
        Inches(0.55),
        Inches(0.85),
        Inches(12),
        Inches(0.7),
        "ПЕРСПЕКТИВЫ: ОТ АНАЛИТИКИ К УПРАВЛЯЕМОМУ РОСТУ",
        size=24,
        bold=True,
        color=C["title"],
    )
    blocks = [
        ("Ближайший этап", C["green"], [
            "• калибровка прогноза по фактическим бронированиям;",
            "• развитие внутреннего Max-бота для сотрудников;",
            "• регулярная проверка рекомендаций и их результатов.",
        ]),
        ("Следующий этап", C["accent"], [
            "• персональные предложения повторным гостям;",
            "• усиление прямых бронирований;",
            "• анализ отзывов и обращений гостей;",
            "• расширение мониторинга конкурентов.",
        ]),
        ("Долгосрочно", C["eyebrow"], [
            "• полуавтоматическое управление тарифами только после подтверждения точности модели;",
            "• масштабирование решения на другие объекты.",
        ]),
    ]
    y = Inches(1.65)
    for head, color, lines in blocks:
        _add_textbox(slide, Inches(0.55), y, Inches(12), Inches(0.35), head, size=18, bold=True, color=color)
        _add_bullets(slide, Inches(0.55), y + Inches(0.35), Inches(12), Inches(1.2), lines, size=15)
        y += Inches(1.75)
    _footer(slide, slide_num)
    _notes(
        slide,
        "Ближайший этап — калибровка прогноза и Max-бот для сотрудников. Дальше — персональные "
        "предложения, прямые брони и отзывы. Полуавтоматическое управление тарифами — только после "
        "подтверждения точности модели на пилоте, не как уже работающая функция.",
    )


def modify_slide5(slide) -> None:
    _replace_exact(slide, "За 2 недели заполняемость в пт-сб ниже будней, а цена держится на прежнем уровне.", "На выходных загрузка ниже будней — система подсказывает, что проверить.")
    _replace_shape_text(slide, "РЕКОМЕНДАЦИИ", "РЕКОМЕНДАЦИЯ: проверить тариф на выходные")
    _notes(
        slide,
        "ИИ читает ваши данные и рынок и говорит по-человечески. Пример: загрузка на выходных "
        "ниже будней — и сразу понятная рекомендация, что проверить. Без формул — вывод и действие.",
    )


def modify_benefit_slide(slide) -> None:
    _replace_exact(slide, "+2–4%", "Цель")
    _replace_exact(
        slide,
        "к среднему чеку и загрузке = рост выручки за сезон",
        "больше выручки за счёт своевременных\nи обоснованных решений по цене, каналам и спросу",
    )
    _replace_exact(
        slide,
        "Меньше рутины · меньше упущенной выручки · больше контроля",
        "Меньше рутины и меньше упущенных возможностей",
    )
    _notes(
        slide,
        "Что это даёт? Десять–пятнадцать минут вместо двух часов каждый день. Цель — больше выручки "
        "за счёт своевременных решений, а не обещание процентов до завершения пилота. Вся доходность — на одном экране.",
    )


def modify_status_slide(slide) -> None:
    mapping = {
        "Система разработана и запущена": "Текущий статус проекта",
        "Готово и работает": "✅ Реализовано",
        "Ядро системы, админка, ИИ-аналитика, конкуренты и тренды, сбор данных, деплой на сервере": (
            "Админка, отчётность, прогноз, рекомендации, конкурентный мониторинг, "
            "события Томска, Word-инструкции и ИИ-аналитика."
        ),
        "Настроено": "⏳ В пилоте",
        "Google-таблицы, почта для отчётов, мессенджер Max, ИИ (YandexGPT)": (
            "Накопление истории TravelLine, проверка точности прогноза, "
            "проверка живых источников событий и цен конкурентов."
        ),
        "Финальный шаг": "➡️ Следующий шаг",
        "Подключение TravelLine и тестовый прогон 1–2 недели перед боевым режимом": (
            "Тестовый прогон 1–2 недели, сверка с TravelLine и настройка рабочих порогов."
        ),
    }
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        txt = shape.text_frame.text.strip()
        if txt in mapping:
            shape.text_frame.text = mapping[txt]
    _notes(
        slide,
        "Реализовано: админка, прогноз, рекомендации, события, Word и ИИ. В пилоте — накопление "
        "истории TravelLine и проверка точности. Следующий шаг — тестовый прогон две недели и сверка с фактом.",
    )


def modify_final_slide(slide) -> None:
    _replace_exact(slide, "ГОТОВЫ ЗАПУСКАТЬ", "ГОТОВЫ ПЕРЕЙТИ К ПИЛОТНОМУ ЗАПУСКУ")
    _replace_exact(slide, "Осталось подключить\nи запустить в работу", "Следующий шаг:\nподключить и проверить данные TravelLine")
    _replace_exact(
        slide,
        "Покажем живую систему, настроим последний источник данных и запустим ежедневные сводки для вашего апарт-отеля.",
        "Запустить тестовый период, сверить прогноз с фактом\nи перейти к регулярной работе.",
    )
    _notes(
        slide,
        "Готовы перейти к пилотному запуску. Следующий шаг — подключить TravelLine, провести "
        "тестовый период, сверить прогноз с фактом и перейти к регулярной работе. Спасибо за внимание!",
    )


def reorder_slides(prs: Presentation, order: list[int]) -> None:
    sldIdLst = prs.slides._sldIdLst
    elements = list(sldIdLst)
    ordered = [elements[i] for i in order]
    for el in list(sldIdLst):
        sldIdLst.remove(el)
    for el in ordered:
        sldIdLst.append(el)


NOTES_BY_SLIDE: dict[int, str] = {
    5: (
        "ИИ читает ваши данные и рынок и говорит по-человечески. Пример: загрузка на выходных "
        "ниже будней — и сразу понятная рекомендация, что проверить. Без формул — вывод и действие."
    ),
    6: (
        "Прогноз — не отчёт за вчера, а инструмент на 7–180 дней вперёд. Система показывает "
        "загрузку, ADR, RevPAR, сценарии и достоверность. На графике — реальные данные из раздела "
        "«Прогноз» с диапазоном неопределённости."
    ),
    7: (
        "Система отслеживает события Томска — конференции, концерты, спорт. Каждое оценивается "
        "по масштабу и вероятности приезда гостей. Это сигнал для проверки цен, а не автоматическая смена тарифа."
    ),
    8: (
        "Рекомендация — не просто цифра, а готовая инструкция: что происходит, что сделать, "
        "как проверить результат и что делать при отклонении. Карточку можно принять, отложить, "
        "назначить сотруднику и выгрузить в Word. Автоматической смены цен в TravelLine нет."
    ),
    9: (
        "Что это даёт? Десять–пятнадцать минут вместо двух часов каждый день. Цель — больше выручки "
        "за счёт своевременных решений, а не обещание процентов до завершения пилота. Вся доходность — на одном экране."
    ),
    10: (
        "Реализовано: админка, прогноз, рекомендации, события, Word и ИИ. В пилоте — накопление "
        "истории TravelLine и проверка точности. Следующий шаг — тестовый прогон две недели и сверка с фактом."
    ),
    11: (
        "Ближайший этап — калибровка прогноза и Max-бот для сотрудников. Дальше — персональные "
        "предложения, прямые брони и отзывы. Полуавтоматическое управление тарифами — только после "
        "подтверждения точности модели на пилоте, не как уже работающая функция."
    ),
    12: (
        "Готовы перейти к пилотному запуску. Следующий шаг — подключить TravelLine, провести "
        "тестовый период, сверить прогноз с фактом и перейти к регулярной работе. Спасибо за внимание!"
    ),
}


def apply_notes_via_com(path: Path) -> None:
    import win32com.client

    app = win32com.client.Dispatch("PowerPoint.Application")
    app.Visible = 1
    try:
        pres = app.Presentations.Open(str(path), WithWindow=False)
        for num, text in NOTES_BY_SLIDE.items():
            slide = pres.Slides(num)
            notes = slide.NotesPage
            if notes.Shapes.Count == 0:
                shape = notes.Shapes.AddShape(1, 36, 36, 648, 504)
                shape.TextFrame.TextRange.Text = text
                shape.Fill.Visible = 0
                shape.Line.Visible = 0
                continue
            try:
                notes.Shapes.Placeholders(2).TextFrame.TextRange.Text = text
                continue
            except Exception:
                pass
            for j in range(1, notes.Shapes.Count + 1):
                shape = notes.Shapes(j)
                if shape.HasTextFrame:
                    shape.TextFrame.TextRange.Text = text
                    break
        pres.Save()
        pres.Close()
    finally:
        app.Quit()


def update_page_numbers(prs: Presentation) -> None:
    for i, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip().isdigit():
                if len(shape.text_frame.text.strip()) <= 2:
                    shape.text_frame.text = str(i)


def main() -> None:
    src = find_source()
    out = output_path()
    shutil.copy2(src, out)
    prs = Presentation(str(out))

    modify_slide5(prs.slides[4])

    build_content_slide(
        prs,
        eyebrow="📈 ПРОГНОЗ",
        title="ПРОГНОЗ ЗАРАНЕЕ, А НЕ ОТЧЁТ ПОСЛЕ ФАКТА",
        body="Система помогает планировать решения на 7, 14, 30 и 180 дней вперёд.",
        bullets=[
            "• прогноз загрузки, ADR, RevPAR и выручки;",
            "• сценарии: базовый, консервативный, оптимистичный;",
            "• учёт текущих броней, темпа новых бронирований и сезонности;",
            "• оценка достоверности прогноза;",
            "• рекомендации по ценам для конкретной даты и категории.",
        ],
        image=ASSETS / "forecast_chart.png",
        slide_num=6,
        notes=(
            "Прогноз — не отчёт за вчера, а инструмент на 7–180 дней вперёд. Система показывает "
            "загрузку, ADR, RevPAR, сценарии и достоверность. На графике — реальные данные из раздела "
            "«Прогноз» с диапазоном неопределённости."
        ),
    )
    build_content_slide(
        prs,
        eyebrow="📅 СОБЫТИЯ",
        title="СОБЫТИЯ ГОРОДА ПОМОГАЮТ УВИДЕТЬ СПРОС ЗАРАНЕЕ",
        body="Система собирает и проверяет события Томска, которые могут влиять на спрос на проживание.",
        bullets=[
            "Конференции · форумы · концерты · спорт · фестивали · праздники",
            "",
            "Событие оценивается по масштабу, длительности, аудитории и вероятности приезда иногородних гостей.",
            "",
            "Подтверждённые события отображаются в прогнозе и помогают заранее проверить цены и доступность.",
        ],
        image=ASSETS / "events_panel.png",
        slide_num=7,
        notes=(
            "Система отслеживает события Томска — конференции, концерты, спорт. Каждое оценивается "
            "по масштабу и вероятности приезда гостей. Это сигнал для проверки цен, а не автоматическая смена тарифа."
        ),
    )
    build_reco_slide(prs, slide_num=8)
    build_perspectives_slide(prs, slide_num=11)

    modify_benefit_slide(prs.slides[5])
    modify_status_slide(prs.slides[6])
    modify_final_slide(prs.slides[7])

    # 0-4 исходные, 8-10 новые (6-8), 5-6 обновлённые (9-10), 11 перспективы, 7 финал
    reorder_slides(prs, [0, 1, 2, 3, 4, 8, 9, 10, 5, 6, 11, 7])
    update_page_numbers(prs)
    prs.save(str(out))
    apply_notes_via_com(out)
    print(f"OK: {out} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
