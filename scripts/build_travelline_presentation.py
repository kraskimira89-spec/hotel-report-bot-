#!/usr/bin/env python3
"""Презентация партнёрского предложения для TravelLine (10 слайдов)."""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUT_DIR = ROOT / "docs" / "presentations" / "travelline"
ASSETS = OUT_DIR / "assets"
OUTPUT = OUT_DIR / "TravelLine_партнёрское_предложение.pptx"

ORG = "Аналитическая платформа для объектов размещения"
CONTACT = "Сергей Богдановский · bogdanchik2@yandex.ru · +7 909 195-04-08"
YEAR = "2026"

DARK = RGBColor(0x0F, 0x17, 0x2A)
DARK_CARD = RGBColor(0x1E, 0x29, 0x3B)
LIGHT_BG = RGBColor(0xF8, 0xFA, 0xFC)
ACCENT = RGBColor(0x38, 0xBD, 0xF8)
ACCENT_DARK = RGBColor(0x1D, 0x4E, 0xD8)
TITLE = RGBColor(0x0F, 0x17, 0x2A)
BODY = RGBColor(0x33, 0x41, 0x55)
MUTED = RGBColor(0x64, 0x74, 0x8B)
WHITE = RGBColor(0xF1, 0xF5, 0xF9)
BORDER = RGBColor(0xE2, 0xE8, 0xF0)

# Крупная типографика для экрана / проектора
F = {
    "eyebrow": 20,
    "title_xl": 48,
    "title_lg": 38,
    "title_md": 32,
    "subtitle": 26,
    "body": 24,
    "body_sm": 22,
    "caption": 18,
    "footer": 13,
    "contact": 22,
}

NOTES: dict[int, str] = {
    1: (
        "Мы не заменяем TravelLine — мы помогаем владельцу быстрее принимать решения на основе "
        "данных, которые уже есть в PMS и каналах продаж. Это партнёрское предложение проверить "
        "ценность вместе, без заявлений о статусе партнёра TravelLine."
    ),
    2: (
        "У клиента TravelLine много данных, но мало времени. Каждый день одни и те же вопросы: "
        "что будет с загрузкой, где цена, какие события, что делать менеджеру. Ответ часто "
        "приходит поздно — после ручного сведения таблиц."
    ),
    3: (
        "Наш слой работает поверх разрешённых данных TravelLine и других источников. "
        "На выходе — не ещё одна таблица, а прогноз, рекомендации и контроль исполнения. "
        "MVP уже проверяется на одном объекте; масштабирование — после пилота."
    ),
    4: (
        "Пример сценария: растёт pickup, есть событие, цена ниже рынка — система предлагает "
        "конкретные шаги и контроль через 24 часа. Автоматической смены тарифов в TravelLine нет — "
        "решение остаётся за менеджером."
    ),
    5: (
        "Мы не обещаем рубли до пилота. Предлагаем измеримую методику: сколько часов уходит на "
        "ручной контроль сейчас и сколько после. Рост дохода измеряем отдельно — не смешиваем "
        "с экономией времени."
    ),
    6: (
        "Для TravelLine это гипотеза: дополнительная ценность данных, регулярность работы с "
        "показателями, add-on без замены PMS. Retention и NPS — только предмет проверки в пилоте, "
        "не заявленный факт."
    ),
    7: (
        "Предлагаем контролируемый пилот с понятным разделением ролей. Коммерческую модель "
        "не фиксируем заранее — сначала подтверждаем спрос и ценность на 1–3 объектах."
    ),
    8: (
        "Вебинар — первый шаг проверки спроса среди клиентов TravelLine. 45 минут, демо на "
        "тестовых данных, опрос готовности платить. Цены в опросе — гипотеза, не прайс."
    ),
    9: (
        "План на 90 дней: согласование API и вебинар, подключение пилота с baseline, анализ "
        "факта и решение о модели партнёрства. Без обещания массового продукта до результатов."
    ),
    10: (
        "Финальный призыв — не продать отчёт, а вместе измерить экономию времени, скорость "
        "решений и готовность платить. Следующий шаг — встреча для согласования вебинара и пилота."
    ),
}


def _run(run, *, size: float, bold: bool = False, color: RGBColor | None = None) -> None:
    run.font.size = Pt(size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color


def _bg(slide, color: RGBColor) -> None:
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def _box(
    slide,
    left,
    top,
    width,
    height,
    text: str,
    *,
    size: float | None = None,
    bold: bool = False,
    color: RGBColor | None = None,
    align=PP_ALIGN.LEFT,
) -> None:
    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = 1.15
    r = p.add_run()
    r.text = text
    _run(r, size=size or F["body"], bold=bold, color=color or BODY)


def _bullets(
    slide,
    left,
    top,
    width,
    height,
    lines: list[str],
    *,
    size: float | None = None,
    color: RGBColor | None = None,
) -> None:
    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame
    tf.word_wrap = True
    sz = size or F["body"]
    col = color or BODY
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(10)
        p.line_spacing = 1.2
        r = p.add_run()
        r.text = line
        _run(r, size=sz, color=col)


def _card(slide, left, top, width, height, fill: RGBColor = LIGHT_BG) -> None:
    s = slide.shapes.add_shape(1, left, top, width, height)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.color.rgb = BORDER


def _footer(slide, n: int, dark: bool = False) -> None:
    c = MUTED if not dark else RGBColor(0x94, 0xA3, 0xB8)
    _box(slide, Inches(0.55), Inches(7.0), Inches(10), Inches(0.25), ORG, size=F["footer"], color=c)
    _box(
        slide,
        Inches(12.0),
        Inches(7.0),
        Inches(0.5),
        Inches(0.25),
        str(n),
        size=F["footer"],
        color=c,
        align=PP_ALIGN.RIGHT,
    )


def _pic(slide, path: Path, left, top, width) -> None:
    if path.exists():
        slide.shapes.add_picture(str(path), left, top, width=width)


def _slide_header(slide, eyebrow: str, title: str, *, dark: bool = False) -> None:
    ec = ACCENT if dark else ACCENT_DARK
    tc = WHITE if dark else TITLE
    _box(slide, Inches(0.55), Inches(0.35), Inches(6), Inches(0.35), eyebrow, size=F["eyebrow"], bold=True, color=ec)
    _box(slide, Inches(0.55), Inches(0.78), Inches(12.2), Inches(0.85), title, size=F["title_lg"], bold=True, color=tc)


def slide_title(prs: Presentation) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, DARK)
    _box(
        s,
        Inches(0.65),
        Inches(0.9),
        Inches(12.0),
        Inches(1.8),
        "От данных TravelLine —\nк ежедневным решениям владельца",
        size=F["title_xl"],
        bold=True,
        color=WHITE,
    )
    _box(
        s,
        Inches(0.65),
        Inches(2.85),
        Inches(12.0),
        Inches(0.7),
        "Партнёрское предложение: аналитика, прогноз и рекомендации",
        size=F["subtitle"],
        color=ACCENT,
    )
    _box(
        s,
        Inches(0.65),
        Inches(3.75),
        Inches(12.0),
        Inches(1.0),
        "Помогаем владельцу быстрее понимать ситуацию,\nдействовать по данным и контролировать результат.",
        size=F["body"],
        color=MUTED,
    )
    _box(s, Inches(0.65), Inches(5.5), Inches(12.0), Inches(0.5), f"{ORG} · {YEAR}", size=F["body_sm"], color=WHITE)
    _box(s, Inches(0.65), Inches(6.15), Inches(12.0), Inches(0.4), CONTACT, size=F["contact"], color=MUTED)
    _footer(s, 1, dark=True)


def slide_problem(prs: Presentation) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, LIGHT_BG)
    _slide_header(s, "ПРОБЛЕМА", "ДАННЫХ МНОГО.\nВРЕМЕНИ НА РЕШЕНИЯ — МАЛО.")
    _bullets(s, Inches(0.55), Inches(1.85), Inches(12.2), Inches(2.4), [
        "• Что будет с загрузкой через неделю и месяц?",
        "• Где цена отстаёт от рынка?",
        "• Какие события могут привести спрос?",
        "• Что должен сделать менеджер сегодня?",
    ], size=F["body"])
    _pic(s, ASSETS / "problem_flow.png", Inches(0.55), Inches(4.35), Inches(12.2))
    _footer(s, 2)


def slide_solution(prs: Presentation) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, DARK)
    _slide_header(s, "РЕШЕНИЕ", "АНАЛИТИЧЕСКИЙ СЛОЙ\nНАД ДАННЫМИ TRAVELLINE", dark=True)
    _pic(s, ASSETS / "solution_stack.png", Inches(0.55), Inches(1.75), Inches(12.2))
    _bullets(s, Inches(0.55), Inches(4.05), Inches(12.2), Inches(2.5), [
        "📊  ежедневная сводка показателей",
        "🔮  прогноз загрузки и выручки",
        "💡  рекомендации и контроль исполнения",
        "📅  события города и рыночные сигналы",
    ], size=F["body"], color=WHITE)
    _box(s, Inches(0.55), Inches(6.35), Inches(12.2), Inches(0.35), "MVP на одном объекте · масштаб — после пилота", size=F["caption"], color=MUTED)
    _footer(s, 3, dark=True)


def slide_owner_value(prs: Presentation) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, LIGHT_BG)
    _slide_header(s, "ЦЕННОСТЬ", "НЕ ОТЧЁТ.\nА ПЛАН ДЕЙСТВИЙ.")
    _card(s, Inches(0.55), Inches(1.85), Inches(12.2), Inches(2.2))
    _box(s, Inches(0.8), Inches(2.05), Inches(11.7), Inches(0.4), "Сигнал:", size=F["body_sm"], bold=True, color=ACCENT_DARK)
    _box(
        s,
        Inches(0.8),
        Inches(2.5),
        Inches(11.7),
        Inches(0.7),
        "Pickup растёт · событие подтверждено · цена ниже рынка",
        size=F["body"],
    )
    _box(s, Inches(0.8), Inches(3.25), Inches(11.7), Inches(0.4), "Действие:", size=F["body_sm"], bold=True, color=ACCENT_DARK)
    _box(
        s,
        Inches(0.8),
        Inches(3.7),
        Inches(11.7),
        Inches(0.6),
        "Проверить цену → обновить тариф → контроль через 24 ч",
        size=F["body"],
    )
    _pic(s, ASSETS / "action_chain.png", Inches(0.55), Inches(4.35), Inches(12.2))
    _footer(s, 4)


def slide_savings(prs: Presentation) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, LIGHT_BG)
    _slide_header(s, "ИЗМЕРЕНИЕ", "ЭКОНОМИЯ ВРЕМЕНИ\nДОЛЖНА БЫТЬ ИЗМЕРИМОЙ")
    _card(s, Inches(0.55), Inches(1.85), Inches(5.9), Inches(2.0))
    _box(
        s,
        Inches(0.8),
        Inches(2.05),
        Inches(5.5),
        Inches(1.6),
        "Часы в месяц =\nежедневный контроль × раб. дни\n+ еженедельный отчёт\n+ сверка данных",
        size=F["body_sm"],
    )
    _card(s, Inches(6.85), Inches(1.85), Inches(5.9), Inches(2.0))
    _box(
        s,
        Inches(7.1),
        Inches(2.05),
        Inches(5.5),
        Inches(1.6),
        "Рубли =\nчасы × стоимость часа\n(ФОТ / 160) × коэфф.",
        size=F["body_sm"],
    )
    _card(s, Inches(0.55), Inches(4.1), Inches(12.2), Inches(1.55), fill=RGBColor(0xEF, 0xF6, 0xFF))
    _box(
        s,
        Inches(0.8),
        Inches(4.3),
        Inches(11.7),
        Inches(1.2),
        "До пилота: ___ ч/мес.   После: ___ ч/мес.\n"
        "Экономия: ___ ч / ___ ₽\n\n"
        "Рост дохода — отдельно, по результатам пилота.",
        size=F["body_sm"],
        color=TITLE,
    )
    _footer(s, 5)


def slide_why_tl(prs: Presentation) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, DARK)
    _slide_header(s, "ЭКОСИСТЕМА", "ЦЕННОСТЬ ДЛЯ TRAVELLINE", dark=True)
    _bullets(s, Inches(0.55), Inches(1.85), Inches(12.2), Inches(4.0), [
        "• Понятная ценность данных для владельца",
        "• Регулярная работа с показателями",
        "• Add-on без замены PMS",
        "• Поэтапное подключение объектов",
        "• Marketplace / совместное предложение",
        "",
        "Retention и NPS — гипотеза для пилота.",
    ], size=F["body"], color=WHITE)
    _footer(s, 6, dark=True)


def slide_partnership(prs: Presentation) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, LIGHT_BG)
    _slide_header(s, "ПИЛОТ", "КОНТРОЛИРУЕМЫЙ СТАРТ")
    rows = [
        ("TravelLine", "API · интеграция · пилотные объекты"),
        ("Разработчик", "аналитика · поддержка · обратная связь"),
        ("Объект", "данные · проверка рекомендаций · оценка"),
    ]
    y = Inches(1.85)
    for role, resp in rows:
        _card(s, Inches(0.55), y, Inches(12.2), Inches(0.95))
        _box(s, Inches(0.8), y + Inches(0.18), Inches(2.8), Inches(0.55), role, size=F["body"], bold=True, color=ACCENT_DARK)
        _box(s, Inches(3.7), y + Inches(0.18), Inches(8.8), Inches(0.55), resp, size=F["body_sm"])
        y += Inches(1.1)
    _box(s, Inches(0.55), Inches(5.25), Inches(12.2), Inches(0.4), "После пилота:", size=F["body"], bold=True, color=TITLE)
    _bullets(s, Inches(0.55), Inches(5.7), Inches(12.2), Inches(1.1), [
        "1. Модуль для клиентов TL   2. Marketplace",
        "3. Совместное внедрение   4. Revenue-share",
    ], size=F["body_sm"])
    _footer(s, 7)


def slide_webinar(prs: Presentation) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, LIGHT_BG)
    _slide_header(s, "СПРОС", "ВЕБИНАР ДЛЯ\nКЛИЕНТОВ TRAVELLINE")
    _card(s, Inches(0.55), Inches(1.85), Inches(12.2), Inches(2.0))
    _box(s, Inches(0.8), Inches(2.05), Inches(11.7), Inches(0.4), "45 минут:", size=F["body"], bold=True, color=ACCENT_DARK)
    _bullets(s, Inches(0.8), Inches(2.5), Inches(11.7), Inches(1.2), [
        "10 мин — проблема владельца",
        "15 мин — демо на тестовых данных",
        "10 мин — прогноз и рекомендации",
        "10 мин — вопросы и заявки в пилот",
    ], size=F["body_sm"])
    _card(s, Inches(0.55), Inches(4.1), Inches(12.2), Inches(1.55))
    _box(s, Inches(0.8), Inches(4.3), Inches(11.7), Inches(0.4), "Опрос:", size=F["body"], bold=True, color=ACCENT_DARK)
    _bullets(s, Inches(0.8), Inches(4.75), Inches(11.7), Inches(0.8), [
        "Формат ценности · готовность платить · заявка в пилот",
        "Цена в опросе — гипотеза, не прайс.",
    ], size=F["body_sm"])
    _footer(s, 8)


def slide_plan90(prs: Presentation) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, DARK)
    _slide_header(s, "ПЛАН", "90 ДНЕЙ ДО ПРОДУКТА", dark=True)
    blocks = [
        ("0–30", "API · вебинар · отбор объектов"),
        ("31–60", "пилот · baseline · обучение"),
        ("61–90", "анализ · спрос · модель партнёрства"),
    ]
    y = Inches(1.85)
    for head, txt in blocks:
        card = s.shapes.add_shape(1, Inches(0.55), y, Inches(12.2), Inches(1.15))
        card.fill.solid()
        card.fill.fore_color.rgb = DARK_CARD
        card.line.color.rgb = BORDER
        _box(s, Inches(0.8), y + Inches(0.22), Inches(1.5), Inches(0.6), head, size=F["title_md"], bold=True, color=ACCENT)
        _box(s, Inches(2.5), y + Inches(0.28), Inches(10.0), Inches(0.6), txt, size=F["body"], color=WHITE)
        y += Inches(1.35)
    _footer(s, 9, dark=True)


def slide_final(prs: Presentation) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, DARK)
    _box(
        s,
        Inches(0.65),
        Inches(1.2),
        Inches(12.0),
        Inches(1.2),
        "ПРОВЕРИМ ЦЕННОСТЬ\nВМЕСТЕ",
        size=F["title_xl"],
        bold=True,
        color=WHITE,
    )
    _box(
        s,
        Inches(0.65),
        Inches(2.65),
        Inches(12.0),
        Inches(1.5),
        "Измерим: экономит ли время,\nпомогает ли действовать быстрее,\nготов ли рынок платить.",
        size=F["body"],
        color=MUTED,
    )
    _box(
        s,
        Inches(0.65),
        Inches(4.4),
        Inches(12.0),
        Inches(0.9),
        "Следующий шаг:\nвстреча · вебинар · пилот",
        size=F["subtitle"],
        bold=True,
        color=ACCENT,
    )
    _box(s, Inches(0.65), Inches(5.8), Inches(12.0), Inches(0.5), CONTACT, size=F["contact"], color=WHITE)
    _footer(s, 10, dark=True)


def apply_notes(path: Path) -> None:
    import win32com.client

    app = win32com.client.Dispatch("PowerPoint.Application")
    app.Visible = 1
    try:
        pres = app.Presentations.Open(str(path.resolve()), WithWindow=False)
        for num, text in NOTES.items():
            slide = pres.Slides(num)
            notes = slide.NotesPage
            if notes.Shapes.Count == 0:
                shape = notes.Shapes.AddShape(1, 36, 36, 648, 504)
                shape.TextFrame.TextRange.Text = text
                shape.Fill.Visible = 0
                shape.Line.Visible = 0
                continue
            try:
                notes.Shapes.Placeholders(2).TextFrame.TextRange.Text = text
            except Exception:
                for j in range(1, notes.Shapes.Count + 1):
                    sh = notes.Shapes(j)
                    if sh.HasTextFrame:
                        sh.TextFrame.TextRange.Text = text
                        break
        pres.Save()
        pres.Close()
    finally:
        app.Quit()


def main() -> None:
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide_title(prs)
    slide_problem(prs)
    slide_solution(prs)
    slide_owner_value(prs)
    slide_savings(prs)
    slide_why_tl(prs)
    slide_partnership(prs)
    slide_webinar(prs)
    slide_plan90(prs)
    slide_final(prs)
    prs.save(str(OUTPUT))
    apply_notes(OUTPUT)
    print(f"OK: {OUTPUT} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
