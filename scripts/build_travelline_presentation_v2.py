#!/usr/bin/env python3
"""TravelLine партнёрское предложение v2 — 12 слайдов.

Стили (типографика, сетка, компоненты) — из figma-design-spec.md.
Цвета — Summer Data Momentum (не из спеки).
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUT_DIR = ROOT / "docs" / "presentations" / "travelline"
SUMMER_DIR = OUT_DIR / "summer"
OUTPUT_PPTX = OUT_DIR / "TravelLine_партнёрское_предложение_v2.pptx"
OUTPUT_PDF = OUT_DIR / "TravelLine_партнёрское_предложение_v2.pdf"
OUTPUT_SUMMER_PPTX = SUMMER_DIR / "TravelLine_Summer_Data_Momentum.pptx"
OUTPUT_SUMMER_PDF = SUMMER_DIR / "TravelLine_Summer_Data_Momentum.pdf"

ORG = "1apart · аналитика для объектов размещения"
CONTACT = "Сергей Богдановский · bogdanchik2@yandex.ru · +7 909 195-04-08"
YEAR = "2026"

# --- сетка 16:9 (figma-design-spec.md, px→in при 1920×1080) ---
# Slide 1920×1080 → 13.333×7.5 in
# Margin 79px → 0.55 in | Content Y 208px → 1.45 in
# Title zone max 151px → 1.05 in | Gap title→content 40px → 0.28 in
# Gap cards 26px → 0.18 in | Footer Y 1018px → 7.05 in
SW = Inches(13.333)
SH = Inches(7.5)
M = Inches(0.55)
CW = SW - M * 2
Y_TITLE = M
H_TITLE = Inches(1.05)
Y_CONTENT = Inches(1.45)  # = title zone end − overlap + gap; контент не выше
GAP_CARD = Inches(0.18)
Y_FOOTER = Inches(7.05)
BAR_W = Inches(6 * 13.333 / 1920)  # 6px accent bar как в Figma
CARD_RADIUS = 0.08  # ~12px на 1920

# --- палитра Summer Data Momentum (цвета НЕ из figma-design-spec) ---
BG_LIGHT = RGBColor(0xF7, 0xFB, 0xFF)
BG_DARK = RGBColor(0x16, 0x78, 0xE8)  # титул/финал: синий градиент-база
TEXT = RGBColor(0x16, 0x25, 0x3A)
ACCENT = RGBColor(0x2D, 0x8C, 0xFF)
ACCENT2 = RGBColor(0x1C, 0xB6, 0xA5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MUTED = RGBColor(0x5A, 0x6E, 0x8C)
CARD = RGBColor(0xFF, 0xFF, 0xFF)
BORDER = RGBColor(0xD8, 0xE8, 0xF5)
SUN = RGBColor(0xFF, 0xC8, 0x4A)
LIGHT_BLUE = RGBColor(0xEA, 0xF5, 0xFF)
LIGHT_MINT = RGBColor(0xE8, 0xFA, 0xF5)
LIGHT_YELLOW = RGBColor(0xFF, 0xF7, 0xDE)
HIGHLIGHT = LIGHT_BLUE
HIGHLIGHT_GREEN = LIGHT_MINT
FONT = "Arial"

# --- типографика (figma-design-spec) ---
F_TITLE_L = 34  # Title/L — обложка, финал
F_TITLE_M = 32  # Title/M — заголовки слайдов
F_SUB = 19      # Subtitle
F_BODY = 16     # Body
F_KPI = 26      # KPI
F_FOOT = 11     # Footer
F_TITLE = F_TITLE_M  # alias для title bar

NOTES: list[str] = [
    # 1
    "Открываем не как замену TravelLine, а как слой решений для владельца на базе уже "
    "существующих данных PMS. Цель встречи — согласовать формат проверки ценности через "
    "вебинар и пилот, без заявлений о статусе официального партнёра.",
    # 2
    "Три боли владельца: данные разрознены, решения запаздывают, нет понятного следующего "
    "шага. Это не критика TravelLine — наоборот, при богатых данных проблема в скорости "
    "интерпретации и действия.",
    # 3
    "Схема: TravelLine остаётся источником истины, мы добавляем прогноз, рекомендации и "
    "контроль. Автоматического изменения тарифов нет — менеджер принимает решение. "
    "Функции пилота помечены как проверяемые.",
    # 4
    "Пример цепочки: сигнал → рекомендация → инструкция → контроль через 24 часа. "
    "Это отличает продукт от статичного отчёта и показывает операционную ценность.",
    # 5
    "Прогноз на несколько горизонтов и календарь событий помогают действовать до того, "
    "как спрос отразится в стандартных отчётах. Уровень уверенности прогноза — часть "
    "честной коммуникации в пилоте.",
    # 6
    "Экономику пилота меряем через часы, не через обещанный рост выручки. Поля для "
    "заполнения — на встрече с объектом. Финансовый эффект — только после baseline и "
    "измерения, без предварительных гарантий.",
    # 7
    "Для TravelLine — дополнительная ценность данных, регулярность, возможный add-on. "
    "Retention и NPS — гипотезы для проверки, не заявленные метрики продукта.",
    # 8
    "Вебинар — низкорисковый шаг: проверить интерес, собрать заявки, понять готовность "
    "платить и приоритеты функций. Таймлайн 45 минут без перегруза.",
    # 9
    "Пилот с разделением ролей: TravelLine — доступ и объекты, разработчик — настройка и "
    "аналитика, объект — обратная связь. Результат — решение о масштабировании на фактах.",
    # 10
    "Три этапа зрелости партнёрства без обещания немедленного marketplace. Модели "
    "монетизации обсуждаются после подтверждённой ценности пилота.",
    # 11
    "90 дней: согласование и вебинар, подключение с baseline, анализ и решение. "
    "Каждый блок имеет измеримый результат, а не только активность.",
    # 12
    "Финальный призыв — совместное измерение: помогает ли решение быстрее действовать. "
    "Следующий шаг — встреча для согласования вебинара и пилота. Контакты на слайде.",
]


def _run(run, *, size: int, bold: bool = False, color: RGBColor | None = None) -> None:
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color


def _bg(slide, color: RGBColor) -> None:
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def _note(slide, text: str) -> None:
    slide.notes_slide.notes_text_frame.text = text


def _rect(slide, left, top, width, height, fill: RGBColor, *, line: RGBColor | None = BORDER) -> None:
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    try:
        sh.adjustments[0] = CARD_RADIUS
    except Exception:
        pass
    if line:
        sh.line.color.rgb = line
    else:
        sh.line.fill.background()
    return sh


def _text(
    slide,
    left,
    top,
    width,
    height,
    text: str,
    *,
    size: int = F_BODY,
    bold: bool = False,
    color: RGBColor = TEXT,
    align=PP_ALIGN.LEFT,
) -> None:
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Pt(4)
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = 1.15
    r = p.add_run()
    r.text = text
    _run(r, size=size, bold=bold, color=color)


def _multiline(
    slide,
    left,
    top,
    width,
    height,
    lines: list[tuple[str, int, bool, RGBColor]],
    *,
    align=PP_ALIGN.LEFT,
) -> None:
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for i, (txt, sz, bold, col) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(6)
        p.line_spacing = 1.12
        r = p.add_run()
        r.text = txt
        _run(r, size=sz, bold=bold, color=col)


def _footer(slide, n: int, *, dark: bool = False) -> None:
    """Slide/Footer — org + номер слайда."""
    c = MUTED if not dark else RGBColor(0xA8, 0xC8, 0xE8)
    _text(slide, M, Y_FOOTER, Inches(8), Inches(0.2), ORG, size=F_FOOT, color=c)
    _text(slide, Inches(12.2), Y_FOOTER, Inches(0.6), Inches(0.2), str(n), size=F_FOOT, color=c, align=PP_ALIGN.RIGHT)


def _slide_title_bar(slide, title: str, *, dark: bool = False) -> None:
    """Slide/TitleBar — Title/M, зона ≤ 1.05 in."""
    tc = WHITE if dark else TEXT
    _text(slide, M, Y_TITLE, CW, H_TITLE, title, size=F_TITLE_M, bold=True, color=tc)


def _hcard(slide, left, top, width, height, title: str, body: str, *, accent: RGBColor = ACCENT) -> None:
    """Slide/Card — rounded + left accent bar + Subtitle + Body."""
    _rect(slide, left, top, width, height, CARD)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, BAR_W, height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()
    pad = Inches(0.14)
    tx = left + Inches(0.18)
    tw = width - Inches(0.22)
    _text(slide, tx, top + pad, tw, Inches(0.35), title, size=F_SUB, bold=True, color=accent)
    _text(slide, tx, top + Inches(0.48), tw, height - Inches(0.55), body, size=F_BODY, color=TEXT)


def _arrow(slide, left, top, width=Inches(0.35)) -> None:
    """Slide/Arrow — chevron между карточками."""
    sh = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, Inches(0.22))
    sh.fill.solid()
    sh.fill.fore_color.rgb = ACCENT2
    sh.line.fill.background()


# --- slides ---


def s01_title(prs: Presentation) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, BG_DARK)
    # декоративная сетка
    for i in range(6):
        ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, M + i * Inches(2.0), Inches(5.8), Inches(1.6), Inches(0.02))
        ln.fill.solid()
        ln.fill.fore_color.rgb = SUN
        ln.line.fill.background()
    _text(
        s, M, Inches(1.0), CW, Inches(1.6),
        "От данных TravelLine —\nк ежедневным решениям владельца",
        size=F_TITLE_L, bold=True, color=WHITE,
    )
    _text(
        s, M, Inches(2.75), CW, Inches(0.9),
        "Партнёрское предложение:\nаналитика, прогноз и рекомендации для объектов размещения",
        size=F_SUB, color=RGBColor(0x8E, 0xC5, 0xFF),
    )
    _text(s, M, Inches(5.6), CW, Inches(0.35), f"{ORG} · {YEAR}", size=F_BODY, color=MUTED)
    _footer(s, 1, dark=True)
    _note(s, NOTES[0])


def s02_problem(prs: Presentation) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, BG_LIGHT)
    _slide_title_bar(s, "Данные есть. Времени на своевременные решения — мало.")
    y = Y_CONTENT
    cards = [
        ("Ручной контроль", "Таблицы, бронирования, цены и отчёты в разных местах."),
        ("Решения с опозданием", "Цена и спрос меняются каждый day."),
        ("Нет следующего действия", "Владелец видит цифры, но не получает понятный план."),
    ]
    # fix typo day -> день
    cards[1] = ("Решения с опозданием", "Цена и спрос меняются каждый день.")
    cw = (CW - GAP_CARD * 2) / 3
    for i, (t, b) in enumerate(cards):
        x = M + i * (cw + GAP_CARD)
        _hcard(s, x, y, cw, Inches(1.55), t, b, accent=ACCENT if i == 0 else ACCENT2)
    _footer(s, 2)
    _note(s, NOTES[1])


def s03_solution(prs: Presentation) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, BG_LIGHT)
    _slide_title_bar(s, "Аналитический слой, который усиливает данные TravelLine")
    y = Y_CONTENT
    cx = M + CW / 2
    steps = [
        "TravelLine + PMS + сигналы рынка",
        "Прогноз и контроль",
        "Рекомендации и действия менеджера",
    ]
    box_w = Inches(5.2)
    box_h = Inches(0.52)
    sy = y
    for i, txt in enumerate(steps):
        _rect(s, cx - box_w / 2, sy, box_w, box_h, HIGHLIGHT, line=ACCENT)
        _text(s, cx - box_w / 2 + Inches(0.12), sy + Inches(0.1), box_w - Inches(0.24), box_h, txt, size=F_BODY, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        sy += box_h + Inches(0.08)
        if i < 2:
            arr = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, cx - Inches(0.12), sy - Inches(0.04), Inches(0.24), Inches(0.18))
            arr.fill.solid()
            arr.fill.fore_color.rgb = ACCENT2
            arr.line.fill.background()
            sy += Inches(0.14)
    y_cards = sy + Inches(0.22)
    feats = ["Ежедневные показатели", "Прогноз загрузки", "Рыночные сигналы", "Контроль действий"]
    cw = (CW - GAP_CARD * 3) / 4
    for i, f in enumerate(feats):
        x = M + i * (cw + GAP_CARD)
        _hcard(s, x, y_cards, cw, Inches(1.2), f, "Проверяется в пилоте", accent=ACCENT2)
    _footer(s, 3)
    _note(s, NOTES[2])


def s04_owner(prs: Presentation) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, BG_LIGHT)
    _slide_title_bar(s, "Не ещё один отчёт. А понятный план действий.")
    y = Y_CONTENT
    steps = [
        ("1. Сигнал", "Растёт pickup, есть событие, цена ниже рынка."),
        ("2. Рекомендация", "Проверить цену конкретной категории."),
        ("3. Инструкция", "Пошаговые действия для менеджера."),
        ("4. Контроль", "Проверить результат через 24 часа."),
    ]
    card_w = Inches(2.55)
    card_h = Inches(1.65)
    arrow_w = Inches(0.28)
    total = 4 * card_w + 3 * arrow_w
    x0 = M + (CW - total) / 2
    x = x0
    for i, (t, b) in enumerate(steps):
        _hcard(s, x, y, card_w, card_h, t, b)
        x += card_w
        if i < 3:
            _arrow(s, x + Inches(0.04), y + Inches(0.7), arrow_w - Inches(0.08))
            x += arrow_w
    _footer(s, 4)
    _note(s, NOTES[3])


def s05_forecast(prs: Presentation) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, BG_LIGHT)
    _slide_title_bar(s, "Спрос можно увидеть раньше, чем он попадёт в отчёт")
    y = Y_CONTENT
    half = (CW - Inches(0.22)) / 2
    _rect(s, M, y, half, Inches(3.35), CARD)
    _text(s, M + Inches(0.18), y + Inches(0.15), half - Inches(0.3), Inches(0.4), "Прогноз на 7 / 14 / 30 / 180 дней", size=F_SUB, bold=True, color=ACCENT)
    _multiline(
        s, M + Inches(0.18), y + Inches(0.55), half - Inches(0.3), Inches(2.5),
        [
            ("• загрузка;", F_BODY, False, TEXT),
            ("• ADR и RevPAR;", F_BODY, False, TEXT),
            ("• выручка;", F_BODY, False, TEXT),
            ("• сценарии и уровень уверенности.", F_BODY, False, TEXT),
        ],
    )
    div_x = M + half + Inches(0.11)
    div = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, div_x, y + Inches(0.2), Inches(0.02), Inches(2.9))
    div.fill.solid()
    div.fill.fore_color.rgb = BORDER
    div.line.fill.background()
    rx = div_x + Inches(0.22)
    _rect(s, rx, y, half, Inches(3.35), CARD)
    _text(s, rx + Inches(0.18), y + Inches(0.15), half - Inches(0.3), Inches(0.4), "События и рынок", size=F_SUB, bold=True, color=ACCENT2)
    _text(
        s, rx + Inches(0.18), y + Inches(0.55), half - Inches(0.3), Inches(0.5),
        "Конференции · концерты · спорт · фестивали",
        size=F_BODY, color=TEXT,
    )
    _text(
        s, rx + Inches(0.18), y + Inches(1.15), half - Inches(0.3), Inches(1.5),
        "Подтверждённые события помогают\nзаранее проверить цены и доступность.",
        size=F_BODY, color=MUTED,
    )
    _footer(s, 5)
    _note(s, NOTES[4])


def s06_economics(prs: Presentation) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, BG_LIGHT)
    _slide_title_bar(s, "Ценность измеряется в сохранённом времени и качестве решений")
    y = Y_CONTENT
    _rect(s, M, y, CW, Inches(1.05), HIGHLIGHT, line=ACCENT)
    _text(
        s, M + Inches(0.2), y + Inches(0.12), CW - Inches(0.4), Inches(0.85),
        "Экономия часов в месяц =\nежедневный ручной контроль + еженедельные отчёты + поиск и сверка данных",
        size=F_SUB, bold=True, color=ACCENT,
    )
    y2 = y + Inches(1.05) + GAP_CARD
    kw = (CW - GAP_CARD * 3) / 4
    kpis = ["До пилота:\n___ ч/мес.", "После пилота:\n___ ч/мес.", "Экономия:\n___ ч/мес.", "Эквивалент:\n___ ₽/мес."]
    for i, k in enumerate(kpis):
        x = M + i * (kw + GAP_CARD)
        _rect(s, x, y2, kw, Inches(1.35), CARD)
        _text(s, x + Inches(0.12), y2 + Inches(0.2), kw - Inches(0.2), Inches(1.0), k, size=F_KPI, bold=True, color=TEXT, align=PP_ALIGN.CENTER)
    _text(
        s, M, y2 + Inches(1.45), CW, Inches(0.45),
        "Финансовый эффект и рост выручки оцениваются отдельно по результатам пилота, без предварительных гарантий.",
        size=F_FOOT, color=MUTED,
    )
    _footer(s, 6)
    _note(s, NOTES[5])


def s07_tl_value(prs: Presentation) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, BG_LIGHT)
    _slide_title_bar(s, "Дополнительная ценность для экосистемы TravelLine")
    y = Y_CONTENT
    items = [
        ("Понятнее данные", "Владелец быстрее понимает ситуацию."),
        ("Больше регулярности", "Показатели превращаются в рабочие решения."),
        ("Новый продуктовый слой", "Возможный add-on или marketplace-интеграция."),
        ("Проверяемая гипотеза", "Ценность подтверждается пилотом, а не обещанием."),
    ]
    cw = (CW - GAP_CARD) / 2
    ch = Inches(1.35)
    for i, (t, b) in enumerate(items):
        col = i % 2
        row = i // 2
        x = M + col * (cw + GAP_CARD)
        yy = y + row * (ch + GAP_CARD)
        _hcard(s, x, yy, cw, ch, t, b, accent=ACCENT if col == 0 else ACCENT2)
    _footer(s, 7)
    _note(s, NOTES[6])


def s08_webinar(prs: Presentation) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, BG_LIGHT)
    _slide_title_bar(s, "Начать можно с вебинара для клиентов TravelLine")
    y = Y_CONTENT
    timeline = [
        ("10 мин", "проблема владельца"),
        ("15 мин", "демонстрация сценариев"),
        ("10 мин", "прогноз и рекомендации"),
        ("10 мин", "вопросы и заявка в пилот"),
    ]
    tw = (CW - GAP_CARD * 3) / 4
    th = Inches(1.15)
    for i, (tm, lbl) in enumerate(timeline):
        x = M + i * (tw + GAP_CARD)
        _rect(s, x, y, tw, th, CARD)
        _text(s, x + Inches(0.1), y + Inches(0.12), tw - Inches(0.2), Inches(0.35), tm, size=F_KPI, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        _text(s, x + Inches(0.1), y + Inches(0.52), tw - Inches(0.2), Inches(0.55), lbl, size=F_BODY, color=TEXT, align=PP_ALIGN.CENTER)
        if i < 3:
            _arrow(s, x + tw + Inches(0.02), y + Inches(0.45), GAP_CARD)
    y2 = y + th + Inches(0.28)
    _rect(s, M, y2, CW, Inches(1.55), HIGHLIGHT_GREEN, line=ACCENT2)
    _text(s, M + Inches(0.18), y2 + Inches(0.12), Inches(1.5), Inches(0.3), "Цели:", size=F_SUB, bold=True, color=ACCENT2)
    _multiline(
        s, M + Inches(0.18), y2 + Inches(0.45), CW - Inches(0.36), Inches(1.0),
        [
            ("• проверить интерес;  • собрать заявки в пилот;  • измерить текущие трудозатраты;", F_BODY, False, TEXT),
            ("• узнать готовность платить;  • определить востребованные функции.", F_BODY, False, TEXT),
        ],
    )
    _footer(s, 8)
    _note(s, NOTES[7])


def s09_pilot(prs: Presentation) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, BG_LIGHT)
    _slide_title_bar(s, "Контролируемый пилот вместо большого риска")
    y = Y_CONTENT
    rows = [
        ("TravelLine", "согласование API и подбор пилотных объектов"),
        ("Разработчик", "настройка, поддержка, аналитика результатов"),
        ("Объект", "доступ к данным и обратная связь менеджера"),
    ]
    rh = Inches(0.72)
    for i, (role, resp) in enumerate(rows):
        yy = y + i * (rh + GAP_CARD)
        _rect(s, M, yy, CW, rh, CARD)
        _text(s, M + Inches(0.15), yy + Inches(0.18), Inches(2.2), rh, role, size=F_SUB, bold=True, color=ACCENT)
        _text(s, M + Inches(2.5), yy + Inches(0.18), CW - Inches(2.7), rh, resp, size=F_BODY, color=TEXT)
    y2 = y + 3 * (rh + GAP_CARD) + Inches(0.05)
    _rect(s, M, y2, CW, Inches(0.85), HIGHLIGHT, line=ACCENT)
    _text(
        s, M + Inches(0.18), y2 + Inches(0.15), CW - Inches(0.36), Inches(0.6),
        "Результат пилота: подтверждённая ценность, спрос, экономика и модель масштабирования.",
        size=F_SUB, bold=True, color=ACCENT,
    )
    _footer(s, 9)
    _note(s, NOTES[8])


def s10_prospects(prs: Presentation) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, BG_LIGHT)
    _slide_title_bar(s, "От пилота — к долгосрочному продукту")
    y = Y_CONTENT
    stages = [
        ("1. Пилот", "Проверка сценариев и спроса."),
        ("2. Продукт", "Дополнительный модуль или marketplace-интеграция."),
        ("3. Масштабирование", "Совместное предложение для сегментов клиентов TravelLine."),
    ]
    sw = (CW - GAP_CARD * 2) / 3
    sh = Inches(1.75)
    for i, (t, b) in enumerate(stages):
        x = M + i * (sw + GAP_CARD)
        _hcard(s, x, y, sw, sh, t, b, accent=[ACCENT, ACCENT2, ACCENT][i])
        if i < 2:
            _arrow(s, x + sw + Inches(0.02), y + Inches(0.75), GAP_CARD)
    _text(
        s, M, y + sh + Inches(0.28), CW, Inches(0.45),
        "Возможные модели: лицензирование · add-on · revenue share · совместное внедрение.",
        size=F_BODY, color=MUTED,
    )
    _footer(s, 10)
    _note(s, NOTES[9])


def s11_plan90(prs: Presentation) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, BG_LIGHT)
    _slide_title_bar(s, "План на 90 дней")
    y = Y_CONTENT
    blocks = [
        ("0–30 дней", "Согласование, вебинар, набор пилотных объектов."),
        ("31–60 дней", "Подключение, baseline показателей, обучение."),
        ("61–90 дней", "Измерение результата, анализ спроса и решение о масштабировании."),
    ]
    bw = (CW - GAP_CARD * 2) / 3
    bh = Inches(2.35)
    for i, (head, txt) in enumerate(blocks):
        x = M + i * (bw + GAP_CARD)
        _rect(s, x, y, bw, bh, CARD)
        num = s.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.15), y + Inches(0.15), Inches(0.45), Inches(0.45))
        num.fill.solid()
        num.fill.fore_color.rgb = ACCENT
        num.line.fill.background()
        _text(s, x + Inches(0.15), y + Inches(0.18), Inches(0.45), Inches(0.4), str(i + 1), size=F_SUB, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        _text(s, x + Inches(0.15), y + Inches(0.72), bw - Inches(0.3), Inches(0.45), head, size=F_SUB, bold=True, color=ACCENT)
        _text(s, x + Inches(0.15), y + Inches(1.2), bw - Inches(0.3), bh - Inches(1.35), txt, size=F_BODY, color=TEXT)
    _footer(s, 11)
    _note(s, NOTES[10])


def s12_final(prs: Presentation) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, BG_DARK)
    _text(
        s, M, Inches(1.15), CW, Inches(0.9),
        "Предлагаем проверить ценность вместе",
        size=F_TITLE_L, bold=True, color=WHITE,
    )
    _text(
        s, M, Inches(2.15), CW, Inches(1.4),
        "Не просто показать отчёт,\nа измерить, помогает ли решение\nвладельцу быстрее действовать по данным.",
        size=F_SUB, color=RGBColor(0xA8, 0xC8, 0xE8),
    )
    _text(
        s, M, Inches(3.85), CW, Inches(0.7),
        "Следующий шаг:\nвстреча для согласования вебинара и пилота.",
        size=F_SUB, bold=True, color=ACCENT2,
    )
    _text(s, M, Inches(5.35), CW, Inches(0.45), CONTACT, size=F_BODY, color=WHITE)
    _footer(s, 12, dark=True)
    _note(s, NOTES[11])


def export_pdf(pptx_path: Path, pdf_path: Path) -> bool:
    try:
        import win32com.client  # type: ignore[import-untyped]
    except ImportError:
        print("PDF: установите pywin32 или экспортируйте вручную из PowerPoint")
        return False
    app = win32com.client.Dispatch("PowerPoint.Application")
    try:
        app.Visible = 1
    except Exception:
        pass
    try:
        pres = app.Presentations.Open(str(pptx_path.resolve()), WithWindow=False)
        pres.SaveAs(str(pdf_path.resolve()), 32)  # ppSaveAsPDF
        pres.Close()
        print(f"OK PDF: {pdf_path}")
        return True
    finally:
        app.Quit()


def main() -> None:
    import shutil

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    SUMMER_DIR.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH
    s01_title(prs)
    s02_problem(prs)
    s03_solution(prs)
    s04_owner(prs)
    s05_forecast(prs)
    s06_economics(prs)
    s07_tl_value(prs)
    s08_webinar(prs)
    s09_pilot(prs)
    s10_prospects(prs)
    s11_plan90(prs)
    s12_final(prs)
    prs.save(str(OUTPUT_PPTX))
    print(f"OK PPTX: {OUTPUT_PPTX} ({len(prs.slides)} slides)")
    shutil.copy2(OUTPUT_PPTX, OUTPUT_SUMMER_PPTX)
    print(f"OK PPTX summer: {OUTPUT_SUMMER_PPTX}")
    if export_pdf(OUTPUT_PPTX, OUTPUT_PDF):
        shutil.copy2(OUTPUT_PDF, OUTPUT_SUMMER_PDF)
        print(f"OK PDF summer: {OUTPUT_SUMMER_PDF}")


if __name__ == "__main__":
    main()
